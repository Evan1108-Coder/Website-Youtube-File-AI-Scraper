from __future__ import annotations

import asyncio
import base64
import csv
import json
from pathlib import Path
import xml.etree.ElementTree as ET
from zipfile import BadZipFile

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from PIL import Image
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError
from pypdf import PdfReader
import pytesseract
from striprtf.striprtf import rtf_to_text

try:
    import pillow_avif  # noqa: F401
except ImportError:  # pragma: no cover
    pillow_avif = None

try:
    import pdfplumber
except ImportError:  # pragma: no cover
    pdfplumber = None

from ai_scraper_bot.models import ExtractedContent, VisualInput
from ai_scraper_bot.services.music_analysis import LocalMusicAnalyzer
from ai_scraper_bot.services.vision import LocalVisionAnalyzer
from ai_scraper_bot.services.video_analysis import LocalVideoAnalyzer
from ai_scraper_bot.services.transcription import TranscriptionService
from ai_scraper_bot.utils.image_loading import load_image_with_fallback


SUPPORTED_FILE_TYPES = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".html",
    ".xml",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".rtf",
    ".png",
    ".avif",
    ".jpg",
    ".jpeg",
    ".mp3",
    ".wav",
    ".m4a",
    ".aac",
    ".flac",
    ".ogg",
    ".mp4",
    ".mov",
}

IMAGE_TYPES = {".png", ".avif", ".jpg", ".jpeg"}
AUDIO_VIDEO_TYPES = {
    ".mp3",
    ".wav",
    ".m4a",
    ".aac",
    ".flac",
    ".ogg",
    ".mp4",
    ".mov",
}
VIDEO_TYPES = {".mp4", ".mov"}
AUDIO_TYPES = AUDIO_VIDEO_TYPES - VIDEO_TYPES


class UnsupportedFileTypeError(RuntimeError):
    pass


class FileParser:
    def __init__(
        self,
        transcription_service: TranscriptionService,
        vision_analyzer: LocalVisionAnalyzer | None = None,
        video_analyzer: LocalVideoAnalyzer | None = None,
        music_analyzer: LocalMusicAnalyzer | None = None,
    ) -> None:
        self.transcription_service = transcription_service
        self.vision_analyzer = vision_analyzer
        self.video_analyzer = video_analyzer
        self.music_analyzer = music_analyzer

    async def parse(self, file_path: Path) -> ExtractedContent:
        extension = file_path.suffix.lower()
        if extension not in SUPPORTED_FILE_TYPES:
            raise UnsupportedFileTypeError(f"Unsupported file type: {extension}")
        await _wait_for_file_ready(file_path)
        if not file_path.exists():
            raise RuntimeError(
                f"The uploaded file was not available at parse time: {file_path}"
            )

        title = file_path.name
        if extension in {".txt", ".md"}:
            return ExtractedContent(
                title,
                await asyncio.to_thread(file_path.read_text, "utf-8", errors="ignore"),
                title,
                metadata={"type": "file"},
            )
        if extension == ".csv":
            return ExtractedContent(title, await asyncio.to_thread(_read_csv, file_path), title, metadata={"type": "file"})
        if extension == ".json":
            return ExtractedContent(title, await asyncio.to_thread(_read_json, file_path), title, metadata={"type": "file"})
        if extension == ".html":
            return ExtractedContent(title, await asyncio.to_thread(_read_html_file, file_path), title, metadata={"type": "file"})
        if extension == ".xml":
            return ExtractedContent(title, await asyncio.to_thread(_read_xml_file, file_path), title, metadata={"type": "file"})
        if extension == ".pdf":
            return ExtractedContent(title, await asyncio.to_thread(_read_pdf, file_path), title, metadata={"type": "file"})
        if extension == ".docx":
            try:
                body = await asyncio.to_thread(_read_docx, file_path)
            except (FileNotFoundError, BadZipFile) as exc:
                raise RuntimeError(_office_file_error_message(file_path, "DOCX", exc)) from exc
            return ExtractedContent(title, body, title, metadata={"type": "file"})
        if extension == ".pptx":
            try:
                body = await asyncio.to_thread(_read_pptx, file_path)
            except (FileNotFoundError, PPTXPackageNotFoundError, BadZipFile) as exc:
                raise RuntimeError(_office_file_error_message(file_path, "PPTX", exc)) from exc
            return ExtractedContent(title, body, title, metadata={"type": "file"})
        if extension == ".xlsx":
            try:
                body = await asyncio.to_thread(_read_xlsx, file_path)
            except (FileNotFoundError, BadZipFile, InvalidFileException) as exc:
                raise RuntimeError(_office_file_error_message(file_path, "XLSX", exc)) from exc
            return ExtractedContent(title, body, title, metadata={"type": "file"})
        if extension == ".rtf":
            return ExtractedContent(title, await asyncio.to_thread(_read_rtf, file_path), title, metadata={"type": "file"})
        if extension in IMAGE_TYPES:
            body, visual_inputs, image_notes = await asyncio.to_thread(_read_image, file_path)
            vision_lines = await self._describe_image(file_path)
            if vision_lines:
                body = "\n\n".join(part for part in (body, vision_lines) if part.strip())
            return ExtractedContent(
                title,
                body,
                title,
                metadata={"type": "image"},
                visual_inputs=visual_inputs,
                issues=image_notes,
                reviewed_media=[f"Image file {file_path.name}"],
            )
        if extension in AUDIO_TYPES:
            duration_minutes = await self.transcription_service.probe_duration_minutes(file_path)
            issues: list[str] = []
            reviewed_media = [f"Audio track from {file_path.name}"]
            transcript = ""
            try:
                transcript = await asyncio.wait_for(
                    self.transcription_service.transcribe_media(file_path, duration_minutes),
                    timeout=_media_transcription_timeout_seconds(duration_minutes),
                )
            except asyncio.TimeoutError:
                issues.append("Audio transcription took too long, so the bot moved on with music analysis and file details only.")
            except Exception as exc:
                issues.append(f"Audio transcription failed: {exc}")

            music_review = await self._analyze_music(file_path, transcript, title)
            body_parts = [
                f"Audio filename: {file_path.name}",
            ]
            if transcript.strip():
                body_parts.extend(["", "Audio transcript:", transcript])
            else:
                body_parts.extend(["", "Audio transcript was not available or was too unclear to recover."])
            if music_review.summary_text:
                body_parts.extend(["", music_review.summary_text])
            metadata = {"type": "media", "media_kind": "audio"}
            metadata.update(music_review.metadata)
            return ExtractedContent(
                title,
                "\n".join(part for part in body_parts if part is not None).strip(),
                title,
                metadata=metadata,
                issues=_dedupe_preserve_order([*issues, *music_review.issues]),
                reviewed_media=_dedupe_preserve_order(
                    [*reviewed_media, *(music_review.reviewed_media or [])]
                ),
            )
        if extension in VIDEO_TYPES:
            duration_minutes = await self.transcription_service.probe_duration_minutes(file_path)
            issues: list[str] = []
            reviewed_media = [f"Video file {file_path.name}"]
            transcript = ""
            try:
                transcript = await asyncio.wait_for(
                    self.transcription_service.transcribe_video_media(file_path, duration_minutes),
                    timeout=_media_transcription_timeout_seconds(duration_minutes),
                )
                reviewed_media.append(f"Extracted audio track from {file_path.name}")
            except asyncio.TimeoutError:
                issues.append("Video transcription took too long, so the bot moved on without a full transcript.")
            except Exception as exc:
                issues.append(_sanitize_media_issue(f"Video transcription failed: {exc}"))

            try:
                video_review = await asyncio.wait_for(
                    self._describe_video(file_path, transcript),
                    timeout=_video_visual_timeout_seconds(duration_minutes),
                )
            except asyncio.TimeoutError:
                video_review = LocalVideoAnalyzerResultFallback()
                video_review.issues.append(
                    "Video visual review took too long, so the bot kept the transcript and file details only."
                )
            except Exception as exc:
                video_review = LocalVideoAnalyzerResultFallback()
                video_review.issues.append(f"Video visual review failed: {exc}")
            music_review = await self._analyze_music(file_path, transcript, title)
            body_parts = [
                f"Video filename: {file_path.name}",
            ]
            if transcript.strip():
                body_parts.extend(["", "Audio transcript:", transcript])
            else:
                body_parts.extend(["", "Audio transcript was not available."])
            if video_review.summary_text:
                body_parts.extend(["", video_review.summary_text])
            elif not issues and not video_review.issues:
                body_parts.extend(["", "No separate visual frame notes were produced for this video."])
            if music_review.summary_text:
                body_parts.extend(["", music_review.summary_text])
            metadata = {"type": "media", "media_kind": "video"}
            metadata.update(music_review.metadata)
            return ExtractedContent(
                title,
                "\n".join(part for part in body_parts if part is not None).strip(),
                title,
                metadata=metadata,
                issues=_dedupe_preserve_order(
                    [_sanitize_media_issue(item) for item in [*issues, *video_review.issues, *music_review.issues]]
                ),
                reviewed_media=_dedupe_preserve_order(
                    [*reviewed_media, *(video_review.reviewed_media or []), *(music_review.reviewed_media or [])]
                ),
                video_interval_history=video_review.interval_history,
            )

        raise UnsupportedFileTypeError(f"Unsupported file type: {extension}")

    async def _describe_image(self, file_path: Path) -> str:
        if self.vision_analyzer is None:
            return ""
        return await self.vision_analyzer.analyze_image_file(file_path)

    async def _describe_video(self, file_path: Path, transcript_text: str):
        if self.video_analyzer is None:
            return LocalVideoAnalyzerResultFallback()
        return await self.video_analyzer.analyze_video_file(file_path, transcript_text=transcript_text)

    async def _analyze_music(self, file_path: Path, transcript_text: str, title: str):
        if self.music_analyzer is None:
            return MusicAnalysisResultFallback()
        return await self.music_analyzer.analyze_media_file(
            file_path,
            transcript_text=transcript_text,
            source_label=title,
        )


class LocalVideoAnalyzerResultFallback:
    def __init__(self) -> None:
        self.summary_text = ""
        self.reviewed_media: list[str] = []
        self.interval_history: list[str] = []
        self.issues: list[str] = []


class MusicAnalysisResultFallback:
    def __init__(self) -> None:
        self.summary_text: str = ""
        self.reviewed_media: list[str] = []
        self.issues: list[str] = []
        self.metadata: dict[str, str] = {}


async def _wait_for_file_ready(file_path: Path, timeout_seconds: float = 3.0) -> None:
    deadline = asyncio.get_running_loop().time() + timeout_seconds
    last_size = -1
    stable_hits = 0
    while asyncio.get_running_loop().time() < deadline:
        if file_path.exists():
            size = file_path.stat().st_size
            if size > 0 and size == last_size:
                stable_hits += 1
                if stable_hits >= 2:
                    return
            else:
                stable_hits = 0
                last_size = size
        await asyncio.sleep(0.15)
    if file_path.exists() and file_path.stat().st_size > 0:
        return


def _office_file_error_message(file_path: Path, kind: str, error: Exception) -> str:
    if not file_path.exists():
        return (
            f"The uploaded {kind} file disappeared before it could be opened: {file_path}. "
            "This usually means the saved upload was not available at parse time."
        )
    return (
        f"The uploaded {kind} file could not be opened as a valid office package. "
        "The file may be corrupted, mislabeled, or not fully uploaded."
    )


def _media_transcription_timeout_seconds(duration_minutes: float | None) -> float:
    if duration_minutes is None:
        return 360.0
    return max(180.0, min(1200.0, duration_minutes * 90.0))


def _video_visual_timeout_seconds(duration_minutes: float | None) -> float:
    if duration_minutes is None:
        return 180.0
    return max(120.0, min(600.0, duration_minutes * 30.0))


def _dedupe_preserve_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for item in items:
        cleaned = item.strip()
        if not cleaned or cleaned in seen:
            continue
        seen.add(cleaned)
        ordered.append(cleaned)
    return ordered


def _sanitize_media_issue(issue: str) -> str:
    lowered = issue.lower()
    if "video transcription failed:" in lowered and (
        "does not contain an audio stream" in lowered
        or "failed to load audio:" in lowered
        or "output file does not contain any stream" in lowered
        or "no audio stream" in lowered
    ):
        return "This video file does not contain a usable audio track, so transcript-based audio analysis was unavailable."
    if "music analysis could not extract an audio sample" in lowered and (
        "does not contain an audio stream" in lowered or "no audio stream" in lowered
    ):
        return "Music analysis could not run because this file does not contain an audio track."
    return issue.strip()


def _read_csv(file_path: Path) -> str:
    rows = []
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _read_json(file_path: Path) -> str:
    return json.dumps(json.loads(file_path.read_text("utf-8", errors="ignore")), indent=2)


def _read_html_file(file_path: Path) -> str:
    soup = BeautifulSoup(file_path.read_text("utf-8", errors="ignore"), "html.parser")
    return soup.get_text("\n", strip=True)


def _read_xml_file(file_path: Path) -> str:
    raw_text = file_path.read_text("utf-8", errors="ignore")
    try:
        root = ET.fromstring(raw_text)
        return "\n".join(element.text.strip() for element in root.iter() if element.text and element.text.strip())
    except ET.ParseError:
        soup = BeautifulSoup(raw_text, "html.parser")
        recovered_text = soup.get_text("\n", strip=True)
        if recovered_text:
            return (
                "Recovered text from malformed XML file. The XML structure was invalid, "
                "but readable content was still extracted:\n\n"
                f"{recovered_text}"
            )
        raise


def _read_pdf(file_path: Path) -> str:
    if pdfplumber is not None:
        extracted_pages = []
        with pdfplumber.open(str(file_path)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                page_lines = [f"Page {index}"]
                text = (page.extract_text() or "").strip()
                if text:
                    page_lines.append(text)
                tables = page.extract_tables() or []
                for table_number, table in enumerate(tables, start=1):
                    page_lines.append(f"Table {table_number} on page {index}")
                    for row in table:
                        values = [
                            str(cell).strip()
                            for cell in row
                            if cell is not None and str(cell).strip()
                        ]
                        if values:
                            page_lines.append(" | ".join(values))
                extracted_pages.append("\n".join(page_lines))
        combined = "\n\n".join(block for block in extracted_pages if block.strip()).strip()
        if combined:
            return combined

    reader = PdfReader(str(file_path))
    return "\n".join(page.extract_text() or "" for page in reader.pages).strip()


def _read_docx(file_path: Path) -> str:
    document = Document(str(file_path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_number, table in enumerate(document.tables, start=1):
        lines.append(f"Table {table_number}")
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_pptx(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    lines = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines.append(f"Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def _read_xlsx(file_path: Path) -> str:
    workbook = load_workbook(str(file_path), data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_rtf(file_path: Path) -> str:
    return rtf_to_text(file_path.read_text("utf-8", errors="ignore"))


def _read_image(file_path: Path) -> tuple[str, list[VisualInput], list[str]]:
    loaded = load_image_with_fallback(file_path)
    image = loaded.image
    width, height = image.size
    ocr_text = pytesseract.image_to_string(image).strip()
    body_parts = [
        f"Image filename: {file_path.name}",
        f"Image size: {width}x{height}",
    ]
    if ocr_text:
        body_parts.append("OCR text:")
        body_parts.append(ocr_text)
    else:
        body_parts.append("OCR text: none detected")

    data_uri = _image_to_data_uri(image, file_path.suffix.lower())
    visual_inputs = []
    if data_uri:
        visual_inputs.append(
            VisualInput(
                kind="image_data",
                value=data_uri,
                label=file_path.name,
            )
        )
    return "\n".join(body_parts), visual_inputs, loaded.notes


def _image_to_data_uri(image: Image.Image, extension: str) -> str:
    import io

    output = io.BytesIO()
    preview = image.copy()
    preview.thumbnail((1024, 1024))
    format_name = "PNG" if extension in {".png", ".avif"} else "JPEG"
    mime_type = "image/png" if format_name == "PNG" else "image/jpeg"
    if format_name == "JPEG" and preview.mode not in {"RGB", "L"}:
        preview = preview.convert("RGB")
    preview.save(output, format=format_name, optimize=True, quality=85)
    encoded = base64.b64encode(output.getvalue()).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"
